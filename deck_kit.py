#!/usr/bin/env python3
"""deck_kit — a McKinsey-grade design system for the OSG Smart Guide deck (16:9).

Single source of visual truth. Every slide module composes from these primitives so
the deck stays consistent: one grid, one palette, one type scale, one footer.

Conventions baked in:
  · action titles (title = the slide's takeaway sentence) + small eyebrow/kicker
  · a thin accent rule under the header
  · restrained palette, generous whitespace, a source line + page number on every page
  · a strict 12-column grid inside a fixed content band (TOP=1.9" .. BOTTOM=6.85")
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# ─────────────────────────── palette (OCBC brand red) ───────────────────────────
# OCBC red is THE accent, used sparingly for emphasis; graphite carries the titles and
# dark grounds; a muted steel + warm ochre are the secondary categoricals; neutral greys
# for text. (Variable names kept stable — NAVY = the dark, AMBER = the red accent.)
NAVY   = RGBColor(0x24, 0x28, 0x2C)   # graphite — titles, dark grounds
NAVY2  = RGBColor(0x33, 0x38, 0x3D)   # lifted graphite for bands / chips on dark
TEAL   = RGBColor(0x4E, 0x6B, 0x84)   # muted steel-blue — secondary categorical
TEAL_D = RGBColor(0x39, 0x52, 0x68)   # darker steel
AMBER  = RGBColor(0xE4, 0x00, 0x2B)   # brand spec red — THE accent, use sparingly
GOLD   = RGBColor(0xB0, 0x81, 0x2C)   # warm ochre — tertiary categorical
INK    = RGBColor(0x20, 0x24, 0x28)   # near-black body text
SLATE  = RGBColor(0x4C, 0x53, 0x59)   # secondary text
GREY   = RGBColor(0x6A, 0x71, 0x77)   # muted / sub text
FAINT  = RGBColor(0x9B, 0xA1, 0xA7)   # captions, page furniture
RULE   = RGBColor(0xE2, 0xE3, 0xE5)   # hairlines / card borders
RULE_D = RGBColor(0xCB, 0xCD, 0xD1)   # slightly darker hairline
CARD   = RGBColor(0xF7, 0xF7, 0xF8)   # card fill (neutral)
CARD2  = RGBColor(0xF1, 0xF1, 0xF3)   # alt card fill
BAND   = RGBColor(0xF4, 0xF1, 0xF2)   # callout band (faint warm grey)
MINT   = RGBColor(0xFB, 0xEE, 0xF0)   # faint red-tint blocks
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PAPER  = RGBColor(0xFC, 0xFC, 0xFC)   # off-white page ground
ONNAVY = RGBColor(0xDA, 0xDC, 0xDF)   # muted text on the dark ground
ONNAVY2= RGBColor(0x9C, 0xA1, 0xA7)   # fainter text on the dark ground

FONT   = "Calibri"
FONT_L = "Calibri Light"

# ─────────────────────────── geometry ───────────────────────────
W, H = Inches(13.333), Inches(7.5)
ML   = Inches(0.75)          # left margin
MR   = Inches(0.75)          # right margin
CW   = Inches(13.333 - 0.75 - 0.75)   # content width = 11.833"
TOP  = Inches(1.90)          # content band top
BOT  = Inches(6.85)          # content band bottom
GUT  = Inches(0.30)          # standard gutter

SOURCE = "Source: OSG Smart Guide · internal analysis"

_state = {"n": 0, "prs": None}


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    _state["prs"] = prs
    _state["n"] = 0
    return prs


def cols(n, gut=GUT, total=CW, x0=ML):
    """Return [(x, w), ...] for n equal columns across `total` with gutter `gut`."""
    g = int(gut) * (n - 1)
    w = Emu(int((int(total) - g) / n))
    return [(Emu(int(x0) + i * (int(w) + int(gut))), w) for i in range(n)]


def spanx(colspec, i, j):
    """x-position and combined width spanning columns i..j (inclusive) of a cols() list."""
    x = colspec[i][0]
    w = Emu(int(colspec[j][0]) + int(colspec[j][1]) - int(x))
    return x, w


# ─────────────────────────── primitives ───────────────────────────
def blank():
    _state["n"] += 1
    return _state["prs"].slides.add_slide(_state["prs"].slide_layouts[6])


def rect(s, x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def hline(s, x, y, w, color=RULE, weight=1.0):
    return rect(s, x, y, w, Pt(weight), color)


def vline(s, x, y, h, color=RULE, weight=1.0):
    return rect(s, x, y, Pt(weight), h, color)


def text(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         sa=6, ls=1.06, wrap=True):
    """paras = list of paragraphs; each paragraph = list of (text, size, color, bold) runs.
    A run may optionally be (text, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0); p.line_spacing = ls
        for run in para:
            t, sz, col, bold = run[0], run[1], run[2], run[3]
            ital = run[4] if len(run) > 4 else False
            r = p.add_run(); r.text = t; f = r.font
            f.size = Pt(sz); f.color.rgb = col; f.bold = bold; f.italic = ital; f.name = FONT
    return tb


def kicker(s, x, y, label, color=AMBER, size=12):
    """Small uppercase eyebrow with tracked caps look (spaces widen it)."""
    return text(s, x, y, Inches(9), Inches(0.3),
                [[("  ".join(list(label.upper())) if False else label.upper(), size, color, True)]], sa=0)


# ─────────────────────────── page scaffold ───────────────────────────
def header(s, eyebrow, action_title, subtitle=None, title_size=24):
    kicker(s, ML, Inches(0.44), eyebrow)
    text(s, ML, Inches(0.74), CW, Inches(0.95), [[(action_title, title_size, NAVY, True)]],
         sa=0, ls=1.02)
    ry = Inches(1.66) if subtitle is None else Inches(1.62)
    hline(s, ML, ry, CW, NAVY, 1.7)
    # short amber tick over the rule (McKinsey accent)
    rect(s, ML, ry, Inches(0.9), Pt(1.7), AMBER)
    if subtitle:
        text(s, ML, Inches(1.74), CW, Inches(0.35), [[(subtitle, 13.5, SLATE, False)]], sa=0)


def footer(s, source=SOURCE):
    hline(s, ML, Inches(7.02), CW, RULE, 0.75)
    text(s, ML, Inches(7.08), Inches(11), Inches(0.3), [[(source, 9, FAINT, False)]], sa=0)
    text(s, Inches(12.35), Inches(7.08), Inches(0.75), Inches(0.3),
         [[("%02d" % _state["n"], 9.5, GREY, True)]], align=PP_ALIGN.RIGHT, sa=0)


def page(eyebrow, title, subtitle=None, source=SOURCE, title_size=24):
    s = blank(); rect(s, 0, 0, W, H, PAPER)
    header(s, eyebrow, title, subtitle, title_size)
    footer(s, source)
    return s


def dark_page():
    """Full-bleed navy ground with a thin amber spine — for cover / dividers / close."""
    s = blank(); rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, W, H, None)  # noop keeps ordering predictable
    rect(s, 0, 0, Inches(0.22), H, AMBER)
    return s


# ─────────────────────────── components ───────────────────────────
def bullets(s, x, y, w, items, size=15, gap=11, head_col=INK, sub_col=GREY,
            marker=AMBER, mk="▪"):
    tb = s.shapes.add_textbox(x, y, w, Inches(0.4)); tf = tb.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        head, sub = it if isinstance(it, tuple) else (it, None)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.08; p.space_before = Pt(0)
        r = p.add_run(); r.text = mk + "  "; r.font.size = Pt(11); r.font.color.rgb = marker; r.font.bold = True
        r2 = p.add_run(); r2.text = head; r2.font.size = Pt(size); r2.font.color.rgb = head_col; r2.font.bold = True
        if sub:
            r3 = p.add_run(); r3.text = "  " + sub; r3.font.size = Pt(size - 0.5); r3.font.color.rgb = sub_col
            r3.font.bold = False
    return tb


def card(s, x, y, w, h, title, body, num=None, accent=TEAL, fill=CARD, title_size=14.5,
         body_size=12, body_col=GREY):
    """Rounded-corner card with a colored top rule and optional big index number."""
    rect(s, x, y, w, h, fill, line=RULE, lw=1.0)
    rect(s, x, y, w, Inches(0.085), accent)
    ty = y + Inches(0.2)
    if num is not None:
        text(s, x + Inches(0.24), y + Inches(0.18), Inches(1.2), Inches(0.4),
             [[(num, 16, accent, True)]], sa=0)
        ty = y + Inches(0.56)
    text(s, x + Inches(0.24), ty, w - Inches(0.48), h - (ty - y) - Inches(0.15),
         [[(title, title_size, NAVY, True)], [(body, body_size, body_col, False)]], sa=5, ls=1.09)


def stat_tile(s, x, y, w, num, label, accent=AMBER, sub=None, num_size=32):
    text(s, x, y, w, Inches(0.7), [[(num, num_size, accent, True)]], align=PP_ALIGN.CENTER, sa=0)
    text(s, x, y + Inches(0.62), w, Inches(0.5), [[(label, 12.5, SLATE, True)]],

         align=PP_ALIGN.CENTER, sa=0)
    if sub:
        text(s, x, y + Inches(0.92), w, Inches(0.4), [[(sub, 9.5, FAINT, False)]],
             align=PP_ALIGN.CENTER, sa=0)


def chip(s, x, y, w, h, label, fill=CARD, line=RULE, txt=NAVY, size=12, bold=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    rect(s, x, y, w, h, fill, line=line, shape=shape)
    text(s, x + Inches(0.14), y, w - Inches(0.28), h, [[(label, size, txt, bold)]],
         anchor=MSO_ANCHOR.MIDDLE, sa=0)


def callout(s, x, y, w, h, lead, body, lead_col=AMBER, fill=BAND, anchor=MSO_ANCHOR.MIDDLE):
    rect(s, x, y, w, h, fill)
    rect(s, x, y, Inches(0.07), h, AMBER)
    runs = []
    if lead:
        runs.append((lead, 14, lead_col, True))
    runs.append((body, 14, NAVY, True))
    text(s, x + Inches(0.28), y, w - Inches(0.5), h, [runs], anchor=anchor, sa=0, ls=1.08)


def chevron_flow(s, x, y, w, h, steps, gap=Inches(0.30)):
    """steps = [(title, body, color)]. Colored blocks joined by chevrons."""
    n = len(steps)
    bw = Emu(int((int(w) - (n - 1) * int(gap)) / n))
    cx = int(x)
    for i, (t, b, col) in enumerate(steps):
        rect(s, Emu(cx), y, bw, h, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, Emu(cx + int(Inches(0.16))), y + Inches(0.18), Emu(int(bw) - int(Inches(0.32))),
             h - Inches(0.3),
             [[(t, 14.5, WHITE, True)], [(b, 10.5, RGBColor(0xEE, 0xEF, 0xF1), False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sa=5, ls=1.06)
        if i < n - 1:
            a = s.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                   Emu(cx + int(bw) + int(gap) // 2 - int(Inches(0.12))),
                                   y + Emu(int(h) // 2 - int(Inches(0.22))),
                                   Inches(0.26), Inches(0.44))
            a.fill.solid(); a.fill.fore_color.rgb = FAINT; a.line.fill.background(); a.shadow.inherit = False
        cx += int(bw) + int(gap)


def column_chart(s, x, y, w, h, series, maxval=None, base_col=TEAL, label_col=SLATE,
                 value_fmt="{:,}"):
    """Simple vertical bar chart drawn with rectangles. series=[(label,value,color?)].
    Baseline at bottom, value labels above bars, category labels below."""
    n = len(series)
    if maxval is None:
        maxval = max(v for _, v, *_ in series)
    plot_h = int(h) - int(Inches(0.85))   # leave room for value + category labels
    plot_top = int(y) + int(Inches(0.30))
    slot = int(w) / n
    bw = int(slot * 0.52)
    base_y = plot_top + plot_h
    hline(s, x, Emu(base_y), w, RULE_D, 1.0)
    for i, item in enumerate(series):
        lab, val = item[0], item[1]
        col = item[2] if len(item) > 2 else base_col
        bh = int(plot_h * (val / maxval)) if maxval else 0
        bx = int(x) + int(i * slot + (slot - bw) / 2)
        by = base_y - bh
        rect(s, Emu(bx), Emu(by), Emu(bw), Emu(bh), col)
        text(s, Emu(bx - int(Inches(0.3))), Emu(by - int(Inches(0.32))),
             Emu(bw + int(Inches(0.6))), Inches(0.3),
             [[(value_fmt.format(val), 13.5, NAVY, True)]], align=PP_ALIGN.CENTER, sa=0)
        text(s, Emu(bx - int(Inches(0.35))), Emu(base_y + int(Inches(0.10))),
             Emu(bw + int(Inches(0.7))), Inches(0.5),
             [[(lab, 11.5, label_col, False)]], align=PP_ALIGN.CENTER, sa=0, ls=1.02)


def legend(s, x, y, entries, size=11.5, gap=Inches(1.9), sw=Inches(0.16)):
    """entries=[(label,color)] laid out horizontally."""
    cx = int(x)
    for lab, col in entries:
        rect(s, Emu(cx), y + Inches(0.03), sw, sw, col)
        text(s, Emu(cx + int(sw) + int(Inches(0.08))), y - Inches(0.02), Inches(1.7), Inches(0.3),
             [[(lab, size, SLATE, False)]], sa=0)
        cx += int(gap)


def divider_number(s, x, y, n, color=AMBER, size=15):
    text(s, x, y, Inches(0.6), Inches(0.4), [[("%02d" % n, size, color, True)]], sa=0)


# ─────── fill-in data placeholders (for figures the client will provide) ───────
def data_slot(s, x, y, w, h, hint="figure", size=12, fill=RGBColor(0xFC, 0xEE, 0xF0)):
    """A dashed, tinted box that visibly reads as 'insert value here'. Use for every
    number that is not yet sourced (payback, CAPEX/OPEX, minutes saved, …)."""
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = AMBER; sp.line.width = Pt(1.1)
    sp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    sp.shadow.inherit = False
    text(s, x, y, w, h, [[(hint, size, AMBER, True, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sa=0)
    return sp


def kpi_slot(s, x, y, w, label, hint="000", sub=None):
    """A headline metric that is still to be provided: a dashed value box + caption."""
    data_slot(s, x + Emu(int((int(w) - int(Inches(1.7))) / 2)), y, Inches(1.7),
              Inches(0.62), hint, size=17)
    text(s, x, y + Inches(0.70), w, Inches(0.4), [[(label, 12.5, SLATE, True)]],
         align=PP_ALIGN.CENTER, sa=0)
    if sub:
        text(s, x, y + Inches(1.00), w, Inches(0.34), [[(sub, 9.5, FAINT, False)]],
             align=PP_ALIGN.CENTER, sa=0)


def lane(s, x, y, w, h, label, tab_fill=NAVY, body_fill=CARD, tab_w=Inches(1.55),
         label_col=WHITE):
    """A horizontal swimlane: a colored label tab on the left + a body area to the right.
    Returns (body_x, body_w) for placing step boxes inside the lane."""
    rect(s, x, y, w, h, body_fill, line=RULE)
    rect(s, x, y, tab_w, h, tab_fill)
    text(s, x + Inches(0.14), y, tab_w - Inches(0.28), h,
         [[(label, 12.5, label_col, True)]], anchor=MSO_ANCHOR.MIDDLE, sa=0, ls=1.05)
    bx = Emu(int(x) + int(tab_w) + int(Inches(0.2)))
    bw = Emu(int(w) - int(tab_w) - int(Inches(0.4)))
    return bx, bw


def step(s, x, y, w, h, title, meta=None, fill=WHITE, line=RULE_D, tcol=NAVY,
         mcol=GREY, tsize=11.5):
    """A small process-step box (title + optional meta line) for use inside a lane."""
    rect(s, x, y, w, h, fill, line=line, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    paras = [[(title, tsize, tcol, True)]]
    if meta:
        paras.append([(meta, 9.5, mcol, False)])
    text(s, x + Inches(0.10), y, w - Inches(0.20), h, paras,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sa=2, ls=1.02)


def arrow(s, x, y, w, h, color=FAINT, shape=MSO_SHAPE.CHEVRON):
    a = s.shapes.add_shape(shape, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background(); a.shadow.inherit = False
    return a


if __name__ == "__main__":
    # smoke test: render one page exercising every component
    import os
    prs = new_deck()
    s = page("Kit smoke test", "Every component renders without overlap or error",
             subtitle="A quick visual regression of the deck_kit primitives")
    c = cols(4)
    for i, (cx, cw) in enumerate(c):
        card(s, cx, TOP + Inches(0.1), cw, Inches(1.6), f"Card {i+1}", "Body copy sample text here.",
             num=str(i + 1), accent=[TEAL, AMBER, GOLD, NAVY][i])
    chevron_flow(s, ML, Inches(3.9), CW, Inches(1.1),
                 [("A", "one", TEAL), ("B", "two", GOLD), ("C", "three", AMBER), ("D", "four", NAVY)])
    callout(s, ML, Inches(5.3), CW, Inches(0.9), "Implication:  ",
            "components compose on a shared grid.")
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_kit_smoke.pptx")
    prs.save(out); print("kit smoke ok →", out, len(prs.slides._sldIdLst), "slide")
