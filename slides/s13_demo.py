#!/usr/bin/env python3
"""Slide 13 — DEMO. The live OSG application screenshot with numbered callouts to the
catalogue, the knowledge graph, the plain-language chat and the source citation.
"""
import os
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import deck_kit as K

IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets", "osg_app.png")


def _badge(s, cx, cy, n, d=Inches(0.36)):
    x = Emu(int(cx) - int(d) // 2)
    y = Emu(int(cy) - int(d) // 2)
    K.rect(s, x, y, d, d, K.AMBER, line=K.WHITE, lw=1.5, shape=MSO_SHAPE.OVAL)
    K.text(s, x, y, d, d, [[(str(n), 12, K.WHITE, True)]], align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE, sa=0)


def render(prs):
    s = K.page("Demo",
               "See it in action: browse the catalogue, explore the knowledge graph, "
               "and ask OSG for a source-cited answer")

    iw, ih = Inches(7.9), Inches(3.71)
    ix, iy = K.ML, Inches(2.55)
    K.rect(s, Emu(int(ix) - int(Inches(0.07))), Emu(int(iy) - int(Inches(0.07))),
           Emu(int(iw) + int(Inches(0.14))), Emu(int(ih) + int(Inches(0.14))),
           K.WHITE, line=K.RULE_D, lw=1.0)
    pic = s.shapes.add_picture(os.path.abspath(IMG), ix, iy, width=iw, height=ih)
    pic.line.color.rgb = K.RULE_D; pic.line.width = Pt(0.75)

    def onimg(fx, fy):
        return int(ix) + int(fx * int(iw)), int(iy) + int(fy * int(ih))

    for n, (fx, fy) in enumerate([(0.10, 0.44), (0.55, 0.55), (0.88, 0.27), (0.80, 0.47)], 1):
        px, py = onimg(fx, fy)
        _badge(s, px, py, n)

    # legend
    lx = Inches(8.95)
    lw = Inches(3.65)
    items = [
        ("Browse the catalogue",
         "2,758 documents in 51 categories, filtered by product or promo."),
        ("Explore the knowledge graph",
         "See how a product links to related and duplicate offers."),
        ("Ask in plain language",
         "OSG replies in seconds, in Bahasa Indonesia or English."),
        ("Every answer cites its source",
         "One click opens the exact document behind the answer."),
    ]
    ly = Inches(2.55)
    for n, (t, d) in enumerate(items, 1):
        _badge(s, int(lx) + int(Inches(0.18)), int(ly) + int(Inches(0.20)), n)
        K.text(s, lx + Inches(0.52), ly, lw - Inches(0.52), Inches(0.9),
               [[(t, 13, K.NAVY, True)], [(d, 11, K.GREY, False)]], sa=2, ls=1.1)
        ly += int(Inches(0.98))
    return s
