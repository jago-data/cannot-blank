#!/usr/bin/env python3
"""Swimlane-flowchart primitives shared by the process-map slides (s07 / s07b).

lanes() draws horizontal actor lanes with rotated left labels; box()/diamond() place
nodes on a grid and return descriptors; connect() draws an elbow arrow between two
nodes with a triangular head and an optional YES/NO pill on the branch.
"""
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
import deck_kit as K

TAB_W = Inches(1.02)
BOX_FILL = RGBColor(0xE9, 0xF4, 0xEF)     # light mint node
BOX_LINE = RGBColor(0xB9, 0xD7, 0xCE)
DIA_FILL = RGBColor(0xEA, 0xF2, 0xF0)
ARROW = RGBColor(0x8B, 0x98, 0x95)


def lanes(s, y, lane_h, defs):
    """defs = [(label, tab_color, body_tint)]. Returns (cols_x0, body_w, [lane_centers])."""
    body_x = Emu(int(K.ML) + int(TAB_W))
    body_w = Emu(int(K.CW) - int(TAB_W))
    centers = []
    for i, (label, tabcol, tint) in enumerate(defs):
        ly = Emu(int(y) + i * int(lane_h))
        K.rect(s, body_x, ly, body_w, lane_h, tint, line=K.RULE)
        K.rect(s, K.ML, ly, TAB_W, lane_h, tabcol)
        lw, lh = lane_h, TAB_W
        cx = int(K.ML) + int(TAB_W) // 2
        cy = int(ly) + int(lane_h) // 2
        tb = s.shapes.add_textbox(Emu(cx - int(lw) // 2), Emu(cy - int(lh) // 2), lw, lh)
        tb.rotation = 270
        tf = tb.text_frame
        tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        f = r.font; f.size = Pt(10.5); f.bold = True; f.color.rgb = K.WHITE; f.name = "Calibri"
        centers.append(cy)
    return int(body_x), int(body_w), centers


def box(s, cx, cy, text, w=Inches(1.74), h=Inches(0.68), fill=BOX_FILL,
        line=BOX_LINE, tcol=K.NAVY, size=9):
    x = Emu(int(cx) - int(w) // 2)
    yy = Emu(int(cy) - int(h) // 2)
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(1.0); sp.shadow.inherit = False
    K.text(s, Emu(int(x) + int(Inches(0.07))), yy, Emu(int(w) - int(Inches(0.14))), h,
           [[(text, size, tcol, False)]], align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE, sa=0, ls=1.02)
    return (int(cx), int(cy), int(w), int(h))


def diamond(s, cx, cy, text, w=Inches(1.78), h=Inches(1.0)):
    x = Emu(int(cx) - int(w) // 2)
    yy = Emu(int(cy) - int(h) // 2)
    sp = s.shapes.add_shape(MSO_SHAPE.DIAMOND, x, yy, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = DIA_FILL
    sp.line.color.rgb = K.TEAL_D; sp.line.width = Pt(1.2); sp.shadow.inherit = False
    K.text(s, Emu(int(x) + int(Inches(0.14))), yy, Emu(int(w) - int(Inches(0.28))), h,
           [[(text, 9, K.NAVY, True)]], align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE, sa=0, ls=1.02)
    return (int(cx), int(cy), int(w), int(h))


def _arrowify(conn, color=ARROW, w=1.5):
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    ln = conn.line._get_or_add_ln()
    for te in ln.findall(qn('a:tailEnd')):
        ln.remove(te)
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))


def _pill(s, mx, my, label, fill):
    w, h = Inches(0.5), Inches(0.28)
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Emu(int(mx) - int(w) // 2), Emu(int(my) - int(h) // 2), w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background(); sp.shadow.inherit = False
    K.text(s, Emu(int(mx) - int(w) // 2), Emu(int(my) - int(h) // 2), w, h,
           [[(label, 8.5, K.WHITE, True)]], align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE, sa=0)


def connect(s, a, b, color=ARROW, branch=None):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    dx, dy = bx - ax, by - ay
    if abs(dx) >= abs(dy):
        x1 = ax + (aw // 2 if dx >= 0 else -aw // 2)
        x2 = bx - (bw // 2 if dx >= 0 else -bw // 2)
        y1, y2 = ay, by
    else:
        y1 = ay + (ah // 2 if dy >= 0 else -ah // 2)
        y2 = by - (bh // 2 if dy >= 0 else -bh // 2)
        x1, x2 = ax, bx
    conn = s.shapes.add_connector(MSO_CONNECTOR.ELBOW, Emu(x1), Emu(y1), Emu(x2), Emu(y2))
    _arrowify(conn, color)
    if branch:
        label, fill = branch
        _pill(s, (x1 + x2) // 2, (y1 + y2) // 2, label, fill)
