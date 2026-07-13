#!/usr/bin/env python3
"""Slide 13a — DEMO (live answer). A real frontliner question ("current credit-card
promos") answered by OSG: a structured, complete reply with amounts, voucher codes,
validity dates and one-click source citations.
"""
import os
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import deck_kit as K

IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets", "osg_promo.png")


def _badge(s, cx, cy, n, d=Inches(0.36)):
    x = Emu(int(cx) - int(d) // 2)
    y = Emu(int(cy) - int(d) // 2)
    K.rect(s, x, y, d, d, K.AMBER, line=K.WHITE, lw=1.5, shape=MSO_SHAPE.OVAL)
    K.text(s, x, y, d, d, [[(str(n), 12, K.WHITE, True)]], align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE, sa=0)


def render(prs):
    s = K.page("Demo",
               "A real question, answered: OSG lists every current card promo with "
               "amounts, codes, dates and sources in seconds")

    iw, ih = Inches(7.9), Inches(3.80)
    ix, iy = K.ML, Inches(2.50)
    K.rect(s, Emu(int(ix) - int(Inches(0.07))), Emu(int(iy) - int(Inches(0.07))),
           Emu(int(iw) + int(Inches(0.14))), Emu(int(ih) + int(Inches(0.14))),
           K.WHITE, line=K.RULE_D, lw=1.0)
    pic = s.shapes.add_picture(os.path.abspath(IMG), ix, iy, width=iw, height=ih)
    pic.line.color.rgb = K.RULE_D; pic.line.width = Pt(0.75)

    def onimg(fx, fy):
        return int(ix) + int(fx * int(iw)), int(iy) + int(fy * int(ih))

    for n, (fx, fy) in enumerate([(0.69, 0.135), (0.40, 0.30), (0.40, 0.66), (0.52, 0.80)], 1):
        px, py = onimg(fx, fy)
        _badge(s, px, py, n)

    # legend
    lx = Inches(8.95)
    lw = Inches(3.65)
    items = [
        ("Ask in plain language",
         "A frontliner types an everyday question, no keywords or product codes needed."),
        ("A complete, structured answer",
         "Five current card promos with amounts, voucher codes and validity dates."),
        ("Every answer cites its source",
         "Zalora, Blibli, Shopee and the bank promo page, one click to verify."),
        ("Guided follow-ups",
         "Suggested next questions open terms and how-to-use in a single tap."),
    ]
    ly = Inches(2.50)
    for n, (t, d) in enumerate(items, 1):
        _badge(s, int(lx) + int(Inches(0.18)), int(ly) + int(Inches(0.20)), n)
        K.text(s, lx + Inches(0.52), ly, lw - Inches(0.52), Inches(0.9),
               [[(t, 13, K.NAVY, True)], [(d, 11, K.GREY, False)]], sa=2, ls=1.1)
        ly += int(Inches(0.98))
    return s
