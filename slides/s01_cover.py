#!/usr/bin/env python3
"""Slide 01 — COVER. Two-column executive hero: left = title, value proposition and
live-proof stats; right = the OSG sign-in product shot; footer = the project team.
"""
import os
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import deck_kit as K

GT = RGBColor(0x2A, 0x74, 0x64)     # gradient start (soft green)
GB = RGBColor(0x08, 0x2C, 0x27)     # gradient end (deep soft green)
SHADOW = RGBColor(0x05, 0x1E, 0x19)
ONGREEN = RGBColor(0xD8, 0xEC, 0xE8)
FRAME = RGBColor(0xEA, 0xF1, 0xEF)
AV_BG = RGBColor(0xE6, 0xEF, 0xEC)
IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets", "osg_login_card.png")


def _grad(s, x, y, w, h, c1, c2, angle=60):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.line.fill.background(); sp.shadow.inherit = False
    sp.fill.gradient()
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    st = sp.fill.gradient_stops
    st[0].color.rgb = c1; st[0].position = 0.0
    st[1].color.rgb = c2; st[1].position = 1.0
    return sp


def _avatar(s, cx, cy, d):
    y = Emu(int(cy) - int(d) // 2)
    K.rect(s, Emu(int(cx) - int(d) // 2), y, d, d, AV_BG, line=K.WHITE, lw=1.5,
           shape=MSO_SHAPE.OVAL)
    hd = Emu(int(int(d) * 0.30))
    K.rect(s, Emu(int(cx) - int(hd) // 2), Emu(int(y) + int(int(d) * 0.16)), hd, hd,
           K.TEAL_D, shape=MSO_SHAPE.OVAL)
    sw = Emu(int(int(d) * 0.60))
    sh = Emu(int(int(d) * 0.46))
    K.rect(s, Emu(int(cx) - int(sw) // 2), Emu(int(y) + int(int(d) * 0.54)), sw, sh,
           K.TEAL_D, shape=MSO_SHAPE.OVAL)


def _stat(s, x, y, num, label):
    K.rect(s, x, Emu(int(y) + int(Inches(0.07))), Inches(0.15), Inches(0.15),
           K.TEAL, shape=MSO_SHAPE.OVAL)
    K.text(s, Emu(int(x) + int(Inches(0.32))), y, Inches(4.9), Inches(0.34),
           [[(num + "   ", 16, K.WHITE, True), (label, 12.5, ONGREEN, False)]],
           anchor=MSO_ANCHOR.MIDDLE, sa=0)


def render(prs):
    s = K.blank()
    _grad(s, 0, 0, K.W, K.H, GT, GB, angle=60)
    K.rect(s, 0, 0, K.W, Inches(0.10), K.TEAL)

    # ── left messaging column ──
    lx = Inches(0.92)
    K.text(s, lx, Inches(0.66), Inches(6), Inches(0.3),
           [[("OSG", 13, K.TEAL, True),
             ("      EXECUTIVE  BRIEFING  ·  BANK  LEADERSHIP", 11.5, ONGREEN, True)]], sa=0)
    K.text(s, Emu(int(lx) - int(Inches(0.02))), Inches(1.08), Inches(5.6), Inches(0.95),
           [[("OSG Smart Guide", 40, K.WHITE, True)]], sa=0)
    K.rect(s, lx, Inches(2.12), Inches(1.05), Inches(0.07), K.TEAL,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    K.text(s, lx, Inches(2.42), Inches(5.35), Inches(1.15),
           [[("A secure AI assistant that gives frontliners instant, source-cited "
              "answers on products, promotions and services, grounded only in the "
              "bank's own content.", 14.5, ONGREEN, False)]], sa=0, ls=1.22)

    # live-proof stats
    K.text(s, lx, Inches(3.78), Inches(5), Inches(0.26),
           [[("LIVE IN PILOT TODAY", 10.5, K.TEAL, True)]], sa=0)
    _stat(s, lx, Inches(4.14), "1,313", "frontliners onboarded")
    _stat(s, lx, Inches(4.66), "12,944", "questions answered")
    _stat(s, lx, Inches(5.18), "326", "products & promos live")

    # ── right product shot (OSG sign-in), 1.66:1 ──
    iw, ih = Inches(6.02), Inches(3.63)
    ix, iy = Inches(6.58), Inches(1.28)
    K.rect(s, Emu(int(ix) + int(Inches(0.06))), Emu(int(iy) + int(Inches(0.11))),
           iw, ih, SHADOW, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    K.rect(s, Emu(int(ix) - int(Inches(0.06))), Emu(int(iy) - int(Inches(0.06))),
           Emu(int(iw) + int(Inches(0.12))), Emu(int(ih) + int(Inches(0.12))),
           FRAME, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    pic = s.shapes.add_picture(os.path.abspath(IMG), ix, iy, width=iw, height=ih)
    pic.line.color.rgb = FRAME; pic.line.width = Pt(0.75)

    # ── project-team footer ──
    K.hline(s, lx, Inches(6.02), Emu(int(K.W) - 2 * int(lx)), RGBColor(0x3A, 0x6C, 0x60), 1.0)
    K.text(s, Inches(0), Inches(6.16), K.W, Inches(0.28),
           [[("PROJECT TEAM", 11, K.TEAL, True)]], align=PP_ALIGN.CENTER, sa=0)
    d = Inches(0.58)
    cxs = (Inches(2.55), Inches(4.75), Inches(6.95), Inches(9.15), Inches(11.35))
    for cx in cxs:
        _avatar(s, cx, Inches(6.86), d)
        K.text(s, Emu(int(cx) - int(Inches(0.7))), Inches(7.22), Inches(1.4), Inches(0.24),
               [[("Name", 9.5, ONGREEN, False)]], align=PP_ALIGN.CENTER, sa=0)
    return s
