#!/usr/bin/env python3
"""OSG — Smart Guide: a McKinsey-style executive deck (16:9), built with python-pptx.

McKinsey conventions used: action titles (each title is the slide's takeaway sentence),
an eyebrow/kicker label, a thin accent rule, a governing-thought executive summary,
restrained palette + whitespace, and a source line on every page.

Reproducible — edit and re-run to rebuild `OSG-Smart-Guide.pptx`.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Restrained, executive palette (OSG brand as the "navy" + minimal accents)
NAVY  = RGBColor(0x0F, 0x3D, 0x38)   # deep green-teal → titles
TEAL  = RGBColor(0x0D, 0x94, 0x88)
ORANGE= RGBColor(0xEA, 0x58, 0x0B)
GOLD  = RGBColor(0xD9, 0x77, 0x06)
INK   = RGBColor(0x25, 0x30, 0x2E)
GREY  = RGBColor(0x60, 0x6B, 0x69)
FAINT = RGBColor(0x94, 0xA3, 0xB8)
RULE  = RGBColor(0xE2, 0xE8, 0xF0)
CARD  = RGBColor(0xF6, 0xF9, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND  = RGBColor(0xEE, 0xF4, 0xF3)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation(); prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
SOURCE = "Source: OSG — Smart Guide (README, backend/wiki.py, knowledge.py, config.yaml)"
_n = [0]

ML = Inches(0.75)          # left margin
CW = Inches(11.83)         # content width


def slide():
    _n[0] += 1
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sa=6):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0); p.line_spacing = 1.05
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t; f = r.font
            f.size = Pt(sz); f.color.rgb = col; f.bold = bold; f.name = "Calibri"
    return tb


def header(s, eyebrow, action_title):
    """McKinsey action-title header: small eyebrow, full-sentence title, thin rule."""
    text(s, ML, Inches(0.42), CW, Inches(0.3), [[(eyebrow.upper(), 11, ORANGE, True)]], sa=0)
    text(s, ML, Inches(0.72), CW, Inches(1.0),
         [[(action_title, 23, NAVY, True)]], sa=0)
    rect(s, ML, Inches(1.62), CW, Pt(1.6), NAVY)


def footer(s):
    text(s, ML, Inches(7.06), Inches(10), Inches(0.3), [[(SOURCE, 8, FAINT, False)]], sa=0)
    text(s, Inches(12.4), Inches(7.06), Inches(0.7), Inches(0.3),
         [[(str(_n[0]), 8, FAINT, False)]], align=PP_ALIGN.RIGHT, sa=0)


def page(eyebrow, title):
    s = slide(); rect(s, 0, 0, W, H, WHITE); header(s, eyebrow, title); footer(s)
    return s


def bullets(s, x, y, w, items, size=14, gap=11, head_col=INK):
    tb = s.shapes.add_textbox(x, y, w, Inches(5)); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        head, sub = it if isinstance(it, tuple) else (it, None)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        r = p.add_run(); r.text = "▪  "; r.font.size = Pt(11); r.font.color.rgb = ORANGE; r.font.bold = True
        r2 = p.add_run(); r2.text = head; r2.font.size = Pt(size); r2.font.color.rgb = head_col; r2.font.bold = True
        if sub:
            r3 = p.add_run(); r3.text = "  " + sub; r3.font.size = Pt(size-0.5); r3.font.color.rgb = GREY
    return tb


def card(s, x, y, w, h, num, title, body, accent=TEAL):
    rect(s, x, y, w, h, CARD, line=RULE)
    rect(s, x, y, w, Inches(0.09), accent)
    if num:
        text(s, x+Inches(0.2), y+Inches(0.16), Inches(0.9), Inches(0.4), [[(num, 15, accent, True)]], sa=0)
    text(s, x+Inches(0.2), y+Inches(0.5 if num else 0.2), w-Inches(0.4), h-Inches(0.6),
         [[(title, 13.5, NAVY, True)], [(body, 11, GREY, False)]], sa=4)


# ══════════════════════════════ 1 — COVER ══════════════════════════════
s = slide(); rect(s, 0, 0, W, H, WHITE)
rect(s, 0, 0, Inches(0.28), H, NAVY)
rect(s, 0, 0, Inches(0.28), Inches(2.2), ORANGE)
text(s, Inches(1.0), Inches(2.35), Inches(11), Inches(2.2),
     [[("OSG — Smart Guide", 42, NAVY, True)],
      [("An internal AI assistant that gives bank staff instant, source-cited answers "
        "on promotions, products and services", 18, GREY, False)]], sa=14)
rect(s, Inches(1.05), Inches(4.65), Inches(3.2), Pt(2.2), ORANGE)
text(s, Inches(1.0), Inches(4.85), Inches(11), Inches(0.6),
     [[("Objective  ·  Challenges  ·  Approach  ·  Demo  ·  Impact  ·  Roadmap", 13, TEAL, True)]], sa=0)
text(s, Inches(1.0), Inches(6.7), Inches(11), Inches(0.4),
     [[("Executive briefing", 11, FAINT, True)]], sa=0)

# ══════════════════════════════ 2 — EXECUTIVE SUMMARY (governing thought) ═════════
s = page("Executive summary",
         "OSG turns a scattered, fast-changing knowledge base into instant, trustworthy answers — "
         "without the fabrication risk of a generic chatbot")
# governing thought band
rect(s, ML, Inches(1.9), CW, Inches(1.15), BAND)
text(s, ML+Inches(0.25), Inches(2.05), CW-Inches(0.5), Inches(0.9),
     [[("Front-line staff ask in plain language; OSG answers only from a curated, self-updating "
        "knowledge base and cites the exact source for every figure — in Bahasa Indonesia or English.",
        15, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE, sa=0)
card(s, ML, Inches(3.35), Inches(3.82), Inches(2.9), "1", "Grounded, never invented",
     "Answers come from source pages only; every claim carries a [n] citation. Expired offers are "
     "flagged as ended, not presented as live.", ORANGE)
card(s, Inches(4.75), Inches(3.35), Inches(3.82), Inches(2.9), "2", "Always current",
     "Drop a new PDF or CMS row and the knowledge hot-reindexes — no restart, no re-training. "
     "Only in-force offers are served.", TEAL)
card(s, Inches(8.68), Inches(3.35), Inches(3.9), Inches(2.9), "3", "Lean & portable",
     "No RAG, no embeddings, no vector DB. A pluggable LLM, a thin stack, and a Docker-free, "
     "root-free deployment.", GOLD)

# ══════════════════════════════ 3 — OBJECTIVE ══════════════════════════════
s = page("Objective",
         "Equip every front-line employee with a trusted, self-updating assistant for the bank's offerings")
bullets(s, ML, Inches(2.0), Inches(6.6), [
    ("Serve the people who serve customers", "relationship managers, tellers and customer service — the “nasabah” is who they help."),
    ("Answer promos, products & services", "what they are, terms, eligibility, and how to apply — bilingual (ID/EN)."),
    ("Make correctness the default", "quote figures exactly; when unsure, say so — never guess."),
    ("Stay current with zero maintenance burden", "the knowledge updates itself as source content changes."),
], size=14.5, gap=15)
# success criteria panel
rect(s, Inches(7.7), Inches(1.95), Inches(4.88), Inches(3.9), CARD, line=RULE)
text(s, Inches(7.95), Inches(2.15), Inches(4.4), Inches(0.4), [[("What “good” looks like", 14, NAVY, True)]], sa=0)
bullets(s, Inches(7.95), Inches(2.7), Inches(4.4), [
    "Every answer is source-cited and verifiable",
    "No expired offer ever shown as active",
    "New content live within seconds, no restart",
    "One model runs selection and synthesis",
    "Runs unprivileged on a locked-down host",
], size=12.5, gap=12)

# ══════════════════════════════ 4 — CHALLENGES ══════════════════════════════
s = page("Challenges",
         "Four structural challenges make manual lookup slow, error-prone and risky")
data = [("A", "Scattered knowledge", "Hundreds of product pages and thousands of time-limited promos, across portals and PDFs.", ORANGE),
        ("B", "Constant change", "Promos start and expire continuously; a stale answer is worse than none.", TEAL),
        ("C", "High cost of error", "A wrong rate or an expired offer quoted to a customer is a compliance & trust risk.", GOLD),
        ("D", "Chatbots hallucinate", "Generic LLMs invent figures and can't cite a source — unacceptable in banking.", NAVY)]
x = ML
for i, (num, t, b, col) in enumerate(data):
    card(s, x, Inches(2.0), Inches(2.86), Inches(3.7), num, t, b, col)
    x += Inches(2.99)
text(s, ML, Inches(5.95), CW, Inches(0.5),
     [[("Implication:  ", 13, ORANGE, True),
       ("staff need answers that are fast AND provably correct AND always up to date — at the same time.",
        13, NAVY, True)]], sa=0)

# ══════════════════════════════ 5 — APPROACH ══════════════════════════════
s = page("Approach",
         "OSG compiles knowledge into a maintained wiki and lets the LLM itself retrieve it — no RAG, no embeddings")
steps = [("Sources", "Cleaned Markdown\nfrom OCBC + PDFs", TEAL),
         ("Wiki", "Category hubs, index,\nhealth, graph", GOLD),
         ("LLM retrieval", "Two stages: pick\ncategory → pick page", ORANGE),
         ("Cited answer", "Grounded, with\n[n] citations", NAVY)]
bw = Inches(2.72); x = ML; y = Inches(2.15)
for i, (t, b, col) in enumerate(steps):
    rect(s, x, y, bw, Inches(1.5), col)
    text(s, x+Inches(0.16), y+Inches(0.2), bw-Inches(0.3), Inches(1.1),
         [[(t, 14.5, WHITE, True)], [(b, 10.5, RGBColor(0xEA,0xF3,0xF1), False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sa=5)
    if i < 3:
        a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x+bw+Emu(15000), y+Inches(0.5), Inches(0.24), Inches(0.5))
        a.fill.solid(); a.fill.fore_color.rgb = FAINT; a.line.fill.background(); a.shadow.inherit=False
    x += bw + Inches(0.32)
text(s, ML, Inches(4.05), CW, Inches(0.4), [[("Why this beats classic RAG", 15, NAVY, True)]], sa=0)
bullets(s, ML, Inches(4.55), CW, [
    ("Knowledge is compiled once and kept current", "not re-derived from raw chunks on every question — cross-references and summaries persist."),
    ("Fully LLM-driven, one model", "the same chat model routes categories then pages; no second model, no embeddings, no vector DB."),
    ("Accuracy guardrails built in", "figures cite the original source; synthesized prose is navigational only; expired offers are flagged, and “nothing relevant” is respected."),
], size=13, gap=9)

# ══════════════════════════════ 6 — DEMO ══════════════════════════════
s = page("Demo",
         "In practice: a staff question returns a concise, source-cited answer in seconds")
# mock chat exchange (based on a real interaction)
rect(s, Inches(7.9), Inches(1.95), Inches(4.68), Inches(0.62), TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(8.05), Inches(2.05), Inches(4.4), Inches(0.45),
     [[("“limit cash collateral loan berapa?”", 12.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE, sa=0)
rect(s, Inches(1.9), Inches(2.75), Inches(6.6), Inches(2.05), WHITE, line=RULE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(2.1), Inches(2.9), Inches(6.25), Inches(1.85),
     [[("Limit Cash Collateral Loan adalah maksimum Rp25 Miliar ", 12.5, INK, False), ("[1]", 10, ORANGE, True), (".", 12.5, INK, False)],
      [("Pinjaman dapat diajukan mulai dari Rp50 Juta hingga Rp25 Miliar, dengan jaminan deposito atau aset cair lainnya ", 12.5, INK, False), ("[1]", 10, ORANGE, True), (".", 12.5, INK, False)],
      [("Mau aku jelaskan syarat atau jangka waktu pinjaman ini?", 12, GREY, False)]], sa=7)
# source chip
rect(s, Inches(1.9), Inches(4.95), Inches(3.3), Inches(0.5), CARD, line=RULE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(2.05), Inches(5.03), Inches(3.1), Inches(0.35),
     [[("[1]  Cash Collateral Loan", 11, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE, sa=0)
text(s, Inches(1.9), Inches(5.6), Inches(6.6), Inches(0.4),
     [[("→ click the source to open the Document View, scoped to that page.", 11, GREY, False)]], sa=0)
# what the demo shows
text(s, Inches(8.9), Inches(2.85), Inches(3.7), Inches(0.4), [[("The demo shows", 14, NAVY, True)]], sa=0)
bullets(s, Inches(8.9), Inches(3.35), Inches(3.7), [
    "Answer-first, numbered, concise",
    "Exact figure quoted from source",
    "[n] citation on every claim",
    "Follow-up offered",
    "One click → the source document",
], size=12, gap=11)

# ══════════════════════════════ 7 — ARCHITECTURE ══════════════════════════════
s = page("How it works",
         "A thin, pluggable stack: React + FastAPI + DuckDB, with a swappable LLM and no vector infrastructure")
boxes = [("Frontend", "React 18 · Vite · Tailwind\nChat · Document View · Admin", TEAL),
         ("Backend API", "FastAPI · Uvicorn\nSSE streaming · LLM-Wiki", ORANGE),
         ("Knowledge store", "DuckDB + Parquet\nMarkdown pages, no DB server", GOLD)]
x = ML
for i, (t, b, col) in enumerate(boxes):
    rect(s, x, Inches(2.15), Inches(3.7), Inches(1.55), CARD, line=RULE)
    rect(s, x, Inches(2.15), Inches(3.7), Inches(0.1), col)
    text(s, x+Inches(0.2), Inches(2.4), Inches(3.35), Inches(1.2),
         [[(t, 14, NAVY, True)], [(b, 11, GREY, False)]], sa=4)
    if i < 2:
        a = s.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, x+Inches(3.73), Inches(2.68), Inches(0.22), Inches(0.42))
        a.fill.solid(); a.fill.fore_color.rgb = FAINT; a.line.fill.background(); a.shadow.inherit=False
    x += Inches(4.02)
rect(s, ML, Inches(4.0), CW, Inches(1.7), BAND)
text(s, ML+Inches(0.25), Inches(4.2), CW-Inches(0.5), Inches(1.4),
     [[("Pluggable LLM — one line in config.yaml", 14, NAVY, True)],
      [("anthropic (Claude Haiku, default)   ·   openai_compatible (remote Qwen / any /v1)   ·   llamacpp (local CPU GGUF)",
        12.5, INK, False)],
      [("Backend :8200  ·  Vite :5174 proxies /api (single origin → no CORS)  ·  ingest via pypdf + watchdog auto-reindex",
        12, GREY, False)]], sa=8)

# ══════════════════════════════ 8 — RESULTS / KB ══════════════════════════════
s = page("Results to date",
         "2,758 real OCBC documents, cleaned and organized into 51 navigable hubs — only in-force offers served")
rect(s, ML, Inches(1.95), CW, Inches(1.45), CARD, line=RULE)
def stat(x, num, lab, col):
    text(s, x, Inches(2.12), Inches(2.7), Inches(1.1),
         [[(num, 30, col, True)], [(lab, 11.5, GREY, False)]], align=PP_ALIGN.CENTER, sa=2)
stat(Inches(0.95), "2,758", "documents ingested", ORANGE)
stat(Inches(3.75), "379", "product pages", TEAL)
stat(Inches(6.5), "2,379", "promo pages", GOLD)
stat(Inches(9.3), "51", "category hubs (37+14)", NAVY)
bullets(s, ML, Inches(3.75), CW, [
    ("Sourced & cleaned from the public OCBC site", "chrome, boilerplate and duplicate “other-product” carousels stripped into focused KB pages."),
    ("Structured for navigation", "products by line (Individu, Korporasi, SME, Syariah…); promos into 14 themes (Dining, Hotel & Travel, Special Event…)."),
    ("Trustworthy by design", "real validity windows gate what's served; a health self-check flags duplicates, stale and orphan pages."),
], size=13.5, gap=11)

# ══════════════════════════════ 9 — IMPACT ══════════════════════════════
s = page("Impact",
         "OSG is designed to cut lookup time and remove fabrication risk while staying maintainable")
card(s, ML, Inches(2.0), Inches(3.82), Inches(3.6), "", "Faster",
     "Seconds to a precise, cited answer instead of searching multiple portals and PDFs — staff stay with the customer.", TEAL)
card(s, Inches(4.75), Inches(2.0), Inches(3.82), Inches(3.6), "", "Safer",
     "Every figure is grounded and cited; expired offers are flagged; customer secrets (OTP/PIN/CVV) are refused by design.", ORANGE)
card(s, Inches(8.68), Inches(2.0), Inches(3.9), Inches(3.6), "", "Cheaper to run & scale",
     "No embeddings/vector DB and a pluggable model keep cost low; self-updating knowledge keeps maintenance near zero.", GOLD)
text(s, ML, Inches(5.85), CW, Inches(0.5),
     [[("Qualitative today; a pilot would quantify time-saved-per-query and answer-accuracy against the current process.",
        11.5, GREY, False)]], sa=0)

# ══════════════════════════════ 10 — ROADMAP ══════════════════════════════
s = page("Roadmap",
         "Next: richer synthesis, deeper facts, and an even friendlier catalog")
phases = [("Now", "Live", "Cleaned KB, category hubs, two-stage LLM retrieval, cited chat, admin dashboard.", TEAL),
          ("Next", "Weeks", "LLM-authored entity & comparison pages with cross-links; Document View catalog revamp.", ORANGE),
          ("Later", "Quarter", "Capture hard facts behind JS tabs (fees/rates/eligibility); retrieval & graph performance; pilot metrics.", GOLD)]
x = ML
for i, (ph, when, body, col) in enumerate(phases):
    rect(s, x, Inches(2.1), Inches(3.86), Inches(3.4), CARD, line=RULE)
    rect(s, x, Inches(2.1), Inches(3.86), Inches(0.6), col)
    text(s, x+Inches(0.22), Inches(2.2), Inches(3.5), Inches(0.42), [[(ph, 15, WHITE, True), ("   "+when, 11, WHITE, False)]], anchor=MSO_ANCHOR.MIDDLE, sa=0)
    text(s, x+Inches(0.22), Inches(2.9), Inches(3.5), Inches(2.4), [[(body, 12.5, GREY, False)]], sa=0)
    x += Inches(4.0)

# ══════════════════════════════ 11 — NEXT STEPS / CLOSE ══════════════════════════════
s = slide(); rect(s, 0, 0, W, H, NAVY)
rect(s, 0, 0, Inches(0.28), H, ORANGE)
text(s, Inches(1.0), Inches(1.4), Inches(11), Inches(0.4), [[("RECOMMENDED NEXT STEPS", 12, GOLD, True)]], sa=0)
bullets_close = [
    ("Run a scoped pilot", "one product line and one branch team; measure time-saved and accuracy vs today."),
    ("Prioritize the hard-facts capture", "fees, rates and eligibility that sit behind JavaScript tabs on the source site."),
    ("Ship the catalog revamp", "structured browse-by-category so staff self-serve beyond chat."),
]
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(3)); tf = tb.text_frame; tf.word_wrap=True
for i, (h_, sub) in enumerate(bullets_close):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(16)
    r = p.add_run(); r.text="▪  "; r.font.size=Pt(13); r.font.color.rgb=GOLD; r.font.bold=True
    r2 = p.add_run(); r2.text=h_; r2.font.size=Pt(17); r2.font.color.rgb=WHITE; r2.font.bold=True
    r3 = p.add_run(); r3.text="  —  "+sub; r3.font.size=Pt(13); r3.font.color.rgb=RGBColor(0xCF,0xE8,0xE4)
rect(s, Inches(1.05), Inches(5.4), Inches(3.0), Pt(2), ORANGE)
text(s, Inches(1.0), Inches(5.6), Inches(11), Inches(0.8),
     [[("OSG — Smart Guide", 22, WHITE, True)],
      [("Grounded · Pluggable · No RAG · Bilingual · Docker-free deploy", 12.5, GOLD, True)]], sa=6)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "OSG-Smart-Guide.pptx")
prs.save(out)
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides)")
